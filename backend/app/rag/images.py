"""Image handling for multi-modal RAG: extract embedded images from slides/PDFs,
caption them with a vision model, and embed the *caption text* rather than the image
itself.

Why caption-then-embed instead of true image embeddings (e.g. CLIP)? Two reasons:
1. It keeps everything in ONE embedding space - a text query ("explain the circuit
   diagram on slide 4") retrieves both text passages and image captions through the
   exact same similarity search, no separate image index or cross-modal matching logic.
2. It's dramatically cheaper/simpler for a prototype at this scale - one small vision
   API call per image at ingestion time, versus running and maintaining a second
   embedding model plus a separate retrieval path that then needs to be merged with
   text results.

This is best-effort: if OPENAI_API_KEY isn't configured or a captioning call fails,
the image is simply skipped rather than failing the whole upload.
"""
import base64
import logging
from pathlib import Path
from typing import Optional

from app.config import settings
from app.observability import trace_ai_call

logger = logging.getLogger("conceptintel")

MAX_IMAGES_PER_FILE = 30  # cap - also a basic zip-bomb / cost guard
MIN_IMAGE_BYTES = 3000  # skip tiny images (bullets, icons, logos) - not worth captioning


def is_captioning_configured() -> bool:
    return bool(settings.OPENAI_API_KEY)


def extract_images_from_pptx(filepath: Path) -> list[bytes]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(filepath)
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    if len(blob) >= MIN_IMAGE_BYTES:
                        images.append(blob)
                except Exception:
                    continue
            if len(images) >= MAX_IMAGES_PER_FILE:
                return images
    return images


def extract_images_from_pdf(filepath: Path) -> list[bytes]:
    import fitz  # PyMuPDF

    images = []
    doc = fitz.open(filepath)
    try:
        for page in doc:
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    blob = base_image["image"]
                    if len(blob) >= MIN_IMAGE_BYTES:
                        images.append(blob)
                except Exception:
                    continue
                if len(images) >= MAX_IMAGES_PER_FILE:
                    return images
    finally:
        doc.close()
    return images


def extract_images(filepath: Path, file_type: str) -> list[bytes]:
    """DOCX image extraction isn't implemented yet (rarer for lecture content than
    PDFs/slides) - noted as future work rather than silently pretending it's covered."""
    if file_type == "pptx":
        return extract_images_from_pptx(filepath)
    if file_type == "pdf":
        return extract_images_from_pdf(filepath)
    return []


@trace_ai_call("image-captioning")
def caption_image(image_bytes: bytes, course_name: str) -> Optional[str]:
    """One vision-model call per image. Returns None (not an exception) on any
    failure or missing config - captioning is a nice-to-have, not required for the
    rest of the pipeline to function."""
    if not is_captioning_configured():
        return None
    try:
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        b64 = base64.b64encode(image_bytes).decode("utf-8")
        response = client.chat.completions.create(
            model=settings.OPENAI_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            f"This image is from a {course_name} lecture slide/document. "
                            "Describe it in 1-3 sentences, focused on the academic concept "
                            "it illustrates (e.g. a diagram, formula, or chart). If it's "
                            "purely decorative (logo, background), say so briefly."
                        ),
                    },
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                ],
            }],
            max_tokens=150,
        )
        return response.choices[0].message.content
    except Exception as e:
        logger.warning("Image captioning failed, skipping this image: %s", e)
        return None
