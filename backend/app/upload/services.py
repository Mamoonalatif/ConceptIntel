import re
import pypdf
import docx
import pptx
import httpx
from pathlib import Path
from app.config import settings


def extract_text_from_file(filepath: Path, file_type: str) -> str:
    """Extracts raw text content from PDF, DOCX, PPTX, or TXT file."""
    file_type = file_type.lower()
    text = ""
    
    if not filepath.exists():
        raise FileNotFoundError(f"File not found at path: {filepath}")

    if file_type == "pdf":
        reader = pypdf.PdfReader(filepath)
        for page in reader.pages:
            text += page.extract_text() or ""
        
    elif file_type == "docx":
        doc = docx.Document(filepath)
        text_list = []
        for paragraph in doc.paragraphs:
            if paragraph.text:
                text_list.append(paragraph.text)
        text = "\n".join(text_list)
        
    elif file_type in ("pptx", "ppt"):
        try:
            prs = pptx.Presentation(filepath)
            text_list = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_list.append(shape.text)
            text = "\n".join(text_list)
        except Exception as e:
            if file_type == "ppt":
                raise ValueError(
                    f"Unable to parse .ppt file. If it is an old binary .ppt format, "
                    f"please convert it to .pptx before uploading. Error: {str(e)}"
                )
            raise ValueError(f"Failed to parse PowerPoint slides: {str(e)}")

        
    elif file_type == "txt":
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
            
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
        
    return clean_extracted_text(text)


def clean_extracted_text(text: str) -> str:
    """Cleans extracted text by removing redundant spacing and formatting."""
    # Replace multiple spaces with a single space
    text = re.sub(r"[ \t]+", " ", text)
    # Replace three or more newlines with double newline
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Remove control characters
    text = "".join(ch for ch in text if ch.isprintable() or ch in "\n\r\t")
    return text.strip()


def chunk_text(text: str, chunk_size: int = 3000, overlap: int = 300) -> list[str]:
    """Splits a long string of text into smaller overlapping chunks."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap
    return chunks


def get_content_type(extension: str) -> str:
    """Returns the MIME content type based on file extension."""
    ext = extension.lower()
    if ext == ".pdf":
        return "application/pdf"
    elif ext == ".docx":
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif ext == ".pptx":
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif ext == ".ppt":
        return "application/vnd.ms-powerpoint"
    elif ext == ".txt":
        return "text/plain"
    return "application/octet-stream"


def _s3_client():
    import boto3
    return boto3.client(
        "s3",
        region_name=settings.AWS_REGION,
        aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
        aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
    )


def _parse_s3_ref(s3_ref: str) -> tuple[str, str]:
    """s3://bucket/key/... -> (bucket, key)."""
    without_scheme = s3_ref[len("s3://"):]
    bucket, _, key = without_scheme.partition("/")
    return bucket, key


def upload_file_to_s3(file_content: bytes, filename: str, content_type: str, course_id: int) -> str:
    """Uploads to a PRIVATE S3 bucket and returns an internal 's3://bucket/key'
    reference (never a public URL - the bucket has no public access, files are only
    ever read back through authenticated boto3 calls in this backend). Raises
    ValueError if AWS credentials/bucket aren't configured."""
    if not (settings.AWS_ACCESS_KEY_ID and settings.AWS_SECRET_ACCESS_KEY and settings.AWS_S3_BUCKET):
        raise ValueError("AWS S3 credentials and bucket are not configured.")

    safe_filename = Path(filename).name
    key = f"course_{course_id}/{safe_filename}"
    client = _s3_client()
    client.put_object(Bucket=settings.AWS_S3_BUCKET, Key=key, Body=file_content, ContentType=content_type)
    return f"s3://{settings.AWS_S3_BUCKET}/{key}"


def download_file_from_s3(s3_ref: str) -> bytes:
    bucket, key = _parse_s3_ref(s3_ref)
    client = _s3_client()
    obj = client.get_object(Bucket=bucket, Key=key)
    return obj["Body"].read()


def delete_file_from_s3(s3_ref: str) -> None:
    bucket, key = _parse_s3_ref(s3_ref)
    client = _s3_client()
    try:
        client.delete_object(Bucket=bucket, Key=key)
    except Exception as e:
        print(f"Warning: Failed to delete S3 object {s3_ref}: {e}")


def _supabase_headers(content_type: str = None) -> dict:
    headers = {
        "Authorization": f"Bearer {settings.SUPABASE_KEY}",
        "apikey": settings.SUPABASE_KEY,
    }
    if content_type:
        headers["Content-Type"] = content_type
    return headers


def _parse_supabase_ref(ref: str) -> tuple[str, str]:
    """supabase://bucket/key/... -> (bucket, key)."""
    without_scheme = ref[len("supabase://"):]
    bucket, _, key = without_scheme.partition("/")
    return bucket, key


def upload_file_to_supabase(file_content: bytes, filename: str, content_type: str, course_id: int) -> str:
    """
    Uploads to Supabase Storage and returns an internal 'supabase://bucket/key'
    reference (NOT a public URL) - the bucket is private, so files are only ever
    read back through authenticated calls (see download_file_from_supabase), same
    approach as the S3 backend. If Supabase settings are missing, raises ValueError.
    """
    if not settings.SUPABASE_URL or not settings.SUPABASE_KEY:
        raise ValueError("Supabase URL and Key are not configured.")

    bucket = settings.SUPABASE_BUCKET
    safe_filename = Path(filename).name
    file_path = f"course_{course_id}/{safe_filename}"

    url = f"{settings.SUPABASE_URL}/storage/v1/object/{bucket}/{file_path}"
    headers = _supabase_headers(content_type)
    headers["x-upsert"] = "true"  # Overwrite if it already exists

    with httpx.Client() as client:
        response = client.post(url, content=file_content, headers=headers, timeout=30.0)

    if response.status_code not in (200, 201):
        raise Exception(f"Supabase storage upload failed with status {response.status_code}: {response.text}")

    return f"supabase://{bucket}/{file_path}"


def download_file_from_supabase(ref: str) -> bytes:
    """Authenticated download from a private Supabase bucket."""
    bucket, key = _parse_supabase_ref(ref)
    url = f"{settings.SUPABASE_URL}/storage/v1/object/{bucket}/{key}"
    with httpx.Client() as client:
        response = client.get(url, headers=_supabase_headers(), timeout=30.0)
    if response.status_code != 200:
        raise Exception(f"Supabase storage download failed with status {response.status_code}: {response.text}")
    return response.content


def delete_file_from_supabase(ref: str) -> None:
    """Deletes a file from Supabase Storage given its 'supabase://bucket/key' reference."""
    if not settings.SUPABASE_URL or not settings.SUPABASE_KEY:
        return
    if not ref.startswith("supabase://"):
        return

    bucket, key = _parse_supabase_ref(ref)
    url = f"{settings.SUPABASE_URL}/storage/v1/object/{bucket}/{key}"

    with httpx.Client() as client:
        response = client.delete(url, headers=_supabase_headers(), timeout=15.0)

    if response.status_code not in (200, 204):
        print(f"Warning: Failed to delete file from Supabase Storage: {response.status_code} - {response.text}")

